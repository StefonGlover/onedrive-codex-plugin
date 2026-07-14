#!/usr/bin/env python3
"""Create genuine rich Office fixtures shared by live beta and host-independent tests."""

import base64
import json
import sys
from io import BytesIO
from pathlib import Path

from PIL import Image
from docx import Document
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.util import Inches


def create_fixtures(root: Path):
    root.mkdir(parents=True, exist_ok=True)
    image_buffer = BytesIO()
    Image.new("RGB", (16, 16), (37, 99, 235)).save(image_buffer, format="PNG")
    image_bytes = image_buffer.getvalue()

    document = Document()
    document.add_heading("Remote Editing Fixture", level=1)
    document.add_paragraph("Word rich operation target")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"; table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "C"; table.cell(1, 1).text = "D"
    document.sections[0].header.paragraphs[0].text = "Original header"
    word_path = root / "rich.docx"
    document.save(word_path)

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["Region", "Revenue"]); sheet.append(["North", 20]); sheet.append(["South", 30])
    revenue_table = Table(displayName="RevenueTable", ref="A1:B3")
    revenue_table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
    sheet.add_table(revenue_table)
    sheet["D1"] = "Code"; sheet["E1"] = "Score"; sheet["D2"] = "B"; sheet["E2"] = 2; sheet["D3"] = "A"; sheet["E3"] = 3
    sheet["F1"] = "Name"; sheet["G1"] = "Value"; sheet["F2"] = "One"; sheet["G2"] = 1; sheet["F3"] = "Two"; sheet["G3"] = 2
    chart = BarChart(); chart.title = "Revenue chart"; chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True); sheet.add_chart(chart, "L1")
    excel_path = root / "rich.xlsx"
    workbook.save(excel_path)

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.75)); box.text = "PowerPoint rich operation target"
    picture = slide.shapes.add_picture(BytesIO(image_bytes), Inches(1), Inches(2), Inches(1), Inches(1))
    table_shape = slide.shapes.add_table(2, 2, Inches(3), Inches(2), Inches(3), Inches(1.5))
    group_a = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(1), Inches(0.5)); group_a.text = "Group A"
    group_b = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(1), Inches(0.5)); group_b.text = "Group B"
    ungroup_a = slide.shapes.add_textbox(Inches(7), Inches(3), Inches(1), Inches(0.5)); ungroup_a.text = "Ungroup A"
    ungroup_b = slide.shapes.add_textbox(Inches(7), Inches(4), Inches(1), Inches(0.5)); ungroup_b.text = "Ungroup B"
    existing_group = slide.shapes.add_group_shape([ungroup_a, ungroup_b])
    powerpoint_path = root / "rich.pptx"
    presentation.save(powerpoint_path)

    return {
        "word": str(word_path),
        "excel": str(excel_path),
        "powerpoint": str(powerpoint_path),
        "imageBase64": base64.b64encode(image_bytes).decode("ascii"),
        "powerpointSelectors": {
            "boxId": str(box.shape_id), "pictureId": str(picture.shape_id), "tableId": str(table_shape.shape_id),
            "groupAId": str(group_a.shape_id), "groupBId": str(group_b.shape_id),
            "existingGroupId": str(existing_group.shape_id), "blankLayoutName": blank_layout.name,
        },
    }


if __name__ == "__main__":
    if len(sys.argv) != 2:
        raise SystemExit("Usage: office-fixture-factory.py <output-directory>")
    print(json.dumps(create_fixtures(Path(sys.argv[1]).resolve())))
