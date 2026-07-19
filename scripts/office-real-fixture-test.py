#!/usr/bin/env python3
"""Interoperability QA using real Office packages and LibreOffice rendering."""

import json
import os
import subprocess
import sys
import tempfile
import base64
import shutil
from io import BytesIO
from contextlib import nullcontext
from pathlib import Path

from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
HELPER = ROOT / "scripts" / "office-openxml.py"
FIXTURE_FACTORY = ROOT / "scripts" / "office-openxml-test.py"
SOFFICE = Path(os.environ.get("SOFFICE") or shutil.which("soffice") or "soffice")
EXPECTED_RICH_OPERATIONS = {
    "word": {"insertImage", "replaceImage", "createContentControl", "deleteContentControl", "createBookmark", "deleteBookmark", "insertTableRow", "deleteTableRow", "insertTableColumn", "deleteTableColumn", "setHeaderFooterText", "setSectionProperties"},
    "excel": {"addWorksheet", "deleteWorksheet", "addTable", "deleteTable", "mergeRange", "unmergeRange", "sortRange", "setAutoFilter", "setHyperlink", "addNote", "deleteNote", "insertImage", "formatChart", "setSheetProtection", "refreshPivot"},
    "powerpoint": {"addSlide", "addImage", "cropImage", "addTable", "insertTableRow", "deleteTableRow", "insertTableColumn", "deleteTableColumn", "setShapeAltText", "setZOrder", "groupShapes", "ungroupShape", "applySlideLayout"},
}


def edit(source, destination, kind, operations):
    result = subprocess.run(
        [sys.executable, str(HELPER)],
        input=json.dumps({"action": "edit", "inputPath": str(source), "outputPath": str(destination), "kind": kind, "operations": operations}),
        text=True,
        capture_output=True,
        check=False,
        env={**os.environ, "PYTHONPYCACHEPREFIX": str(destination.parent / "pycache")},
    )
    payload = json.loads(result.stdout or "{}")
    if result.returncode or not payload.get("ok"):
        raise RuntimeError(payload.get("error") or result.stderr)
    return payload["value"]


def main():
    checks = {}
    kept_root = os.environ.get("OFFICE_REAL_FIXTURE_DIR")
    context = nullcontext(kept_root) if kept_root else tempfile.TemporaryDirectory(prefix="onedrive-office-real-")
    with context as directory:
        root = Path(directory)
        root.mkdir(parents=True, exist_ok=True)
        output = root / "rendered"
        output.mkdir()

        emitted_root = root / "consolidated-fixtures"
        emitted = subprocess.run(
            [sys.executable, str(FIXTURE_FACTORY), f"--emit-fixtures={emitted_root}"],
            capture_output=True,
            text=True,
            check=False,
            env={**os.environ, "PYTHONPYCACHEPREFIX": str(root / "pycache")},
        )
        if emitted.returncode != 0:
            raise RuntimeError(emitted.stderr or emitted.stdout)
        emitted_paths = {
            "word": emitted_root / "sample.docx",
            "excel": emitted_root / "sample.xlsx",
            "powerpoint": emitted_root / "sample.pptx",
        }
        Document(emitted_paths["word"])
        load_workbook(emitted_paths["excel"], data_only=False)
        Presentation(emitted_paths["powerpoint"])
        emitted_rendered = {}
        for kind, path in emitted_paths.items():
            destination = output / ("consolidated-" + kind)
            destination.mkdir()
            profile = root / ("profile-consolidated-" + kind)
            conversion = subprocess.run([str(SOFFICE), "--headless", f"-env:UserInstallation={profile.as_uri()}", "--convert-to", "pdf", "--outdir", str(destination), str(path)], cwd=root, capture_output=True, text=True, check=False)
            pdf = destination / "sample.pdf"
            diagnostic = (conversion.stdout + "\n" + conversion.stderr).lower()
            emitted_rendered[kind] = {
                "exitCode": conversion.returncode,
                "bytes": pdf.stat().st_size if pdf.exists() else 0,
                "cleanOpen": conversion.returncode == 0 and not any(term in diagnostic for term in ("repair", "corrupt", "fatal error")),
                "stderr": conversion.stderr[-500:],
            }

        document = Document()
        document.add_heading("Real Word Fixture", level=1)
        document.add_paragraph("Hello from Word")
        document.add_paragraph("Content control target")
        word_table = document.add_table(rows=2, cols=2)
        word_table.cell(0, 0).text = "A"; word_table.cell(0, 1).text = "B"; word_table.cell(1, 0).text = "C"; word_table.cell(1, 1).text = "D"
        document.sections[0].header.paragraphs[0].text = "Original header"
        source_docx = root / "real.docx"
        document.save(source_docx)
        edited_docx = root / "edited-word.docx"
        word = edit(source_docx, edited_docx, "word", [{"type": "replaceText", "find": "Hello from Word", "replace": "Edited in OneDrive"}, {"type": "insertParagraph", "text": "Inserted safely"}])
        rich_word_docx = root / "rich-word.docx"
        image_b64 = base64.b64encode(Image.new("RGB", (8, 8), (37, 99, 235)).tobytes()).decode("ascii")
        image_buffer = BytesIO(); Image.new("RGB", (8, 8), (37, 99, 235)).save(image_buffer, format="PNG"); image_b64 = base64.b64encode(image_buffer.getvalue()).decode("ascii")
        rich_word_operations = [
            {"type": "insertImage", "paragraphIndex": 0, "base64": image_b64, "contentType": "image/png", "width": 457200, "height": 457200, "altText": "Blue square"},
            {"type": "replaceImage", "imageIndex": 0, "base64": image_b64, "contentType": "image/png"},
            {"type": "createContentControl", "paragraphIndex": 2, "tag": "codex", "title": "Codex target", "text": "Controlled"},
            {"type": "deleteContentControl", "contentControlIndex": 0, "preserveContent": True},
            {"type": "createBookmark", "paragraphIndex": 0, "name": "CodexBookmark"},
            {"type": "deleteBookmark", "name": "CodexBookmark"},
            {"type": "insertTableRow", "tableIndex": 0, "rowIndex": 1, "values": ["E", "F"]},
            {"type": "insertTableColumn", "tableIndex": 0, "columnIndex": 1, "values": ["X", "Y", "Z"]},
            {"type": "deleteTableRow", "tableIndex": 0, "rowIndex": 2},
            {"type": "deleteTableColumn", "tableIndex": 0, "columnIndex": 2},
            {"type": "setHeaderFooterText", "part": "word/header1.xml", "text": "Updated header"},
            {"type": "setSectionProperties", "sectionIndex": 0, "orientation": "landscape", "pageWidth": 15840, "pageHeight": 12240, "marginLeft": 720, "marginRight": 720},
        ]
        rich_word = edit(source_docx, rich_word_docx, "word", rich_word_operations)

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Data"
        sheet.append(["Region", "Revenue"])
        sheet.append(["North", 20])
        sheet.append(["South", 30])
        table = Table(displayName="RevenueTable", ref="A1:B3")
        table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
        sheet.add_table(table)
        sheet["D1"] = "Code"; sheet["E1"] = "Score"; sheet["D2"] = "B"; sheet["E2"] = 2; sheet["D3"] = "A"; sheet["E3"] = 3
        sheet["F1"] = "Name"; sheet["G1"] = "Value"; sheet["F2"] = "One"; sheet["G2"] = 1; sheet["F3"] = "Two"; sheet["G3"] = 2
        chart = BarChart(); chart.title = "Revenue chart"; chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True); sheet.add_chart(chart, "L1")
        source_xlsx = root / "real.xlsx"
        workbook.save(source_xlsx)
        edited_xlsx = root / "edited-excel.xlsx"
        excel = edit(source_xlsx, edited_xlsx, "excel", [
            {"type": "addTableRow", "table": "RevenueTable", "values": [["West", 40]]},
            {"type": "deleteTableRow", "table": "RevenueTable", "index": 0},
            {"type": "setTableTotals", "table": "RevenueTable", "enabled": True, "columns": [{"column": "Region", "label": "Total"}, {"column": "Revenue", "function": "sum"}]},
            {"type": "createChart", "sheet": "Data", "sourceData": "A1:B4", "chartType": "ColumnClustered", "name": "Revenue chart", "titleText": "Revenue by region", "left": 20, "top": 30, "width": 420, "height": 240},
            {"type": "updateChart", "sheet": "Data", "chart": "Revenue chart", "chartType": "Line", "sourceData": "A1:B4", "name": "Revenue trend", "titleText": "Revenue trend", "left": 40, "top": 50, "width": 400, "height": 220},
        ])
        rich_xlsx = root / "rich-excel.xlsx"
        rich_excel_operations = [
            {"type": "addWorksheet", "name": "Scratch"}, {"type": "deleteWorksheet", "sheet": "Scratch"},
            {"type": "addTable", "sheet": "Data", "address": "F1:G3", "name": "CodexTable"}, {"type": "deleteTable", "table": "CodexTable", "preserveData": True},
            {"type": "mergeRange", "sheet": "Data", "address": "H1:I1"}, {"type": "unmergeRange", "sheet": "Data", "address": "H1:I1"},
            {"type": "sortRange", "sheet": "Data", "address": "D2:E3", "keys": [{"column": 1, "descending": True}], "hasHeaders": False},
            {"type": "setAutoFilter", "sheet": "Data", "address": "D1:E3"},
            {"type": "setHyperlink", "sheet": "Data", "address": "J1", "url": "https://openai.com", "display": "OpenAI"},
            {"type": "addNote", "sheet": "Data", "address": "K1", "text": "Temporary note", "author": "Codex"}, {"type": "deleteNote", "sheet": "Data", "address": "K1"},
            {"type": "insertImage", "sheet": "Data", "fromAddress": "N1", "base64": image_b64, "contentType": "image/png"},
            {"type": "formatChart", "sheet": "Data", "chart": "0", "titleText": "Formatted chart", "legendPosition": "bottom", "style": 10},
            {"type": "setSheetProtection", "sheet": "Data", "enabled": True, "allowSelectUnlockedCells": True},
            {"type": "refreshPivot"},
        ]
        rich_excel = edit(source_xlsx, rich_xlsx, "excel", rich_excel_operations)
        mixed_excel_path = root / "mixed-order-excel.xlsx"
        mixed_excel = edit(source_xlsx, mixed_excel_path, "excel", [
            {"type": "addWorksheet", "name": "Interleaved"},
            {"type": "setCell", "sheet": "Interleaved", "address": "A1", "value": "ordered"},
            {"type": "setHyperlink", "sheet": "Interleaved", "address": "A2", "url": "https://openai.com", "display": "OpenAI"},
        ])

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text = "Hello from PowerPoint"
        disposable = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(0.5))
        disposable.text = "Delete me"
        source_png_buffer = BytesIO()
        Image.new("RGB", (8, 8), (37, 99, 235)).save(source_png_buffer, format="PNG")
        source_png = source_png_buffer.getvalue()
        picture = slide.shapes.add_picture(BytesIO(source_png), Inches(1), Inches(3), Inches(1), Inches(1))
        table_shape = slide.shapes.add_table(2, 2, Inches(3), Inches(3), Inches(3), Inches(1.5))
        group_a = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(1), Inches(.5)); group_a.text = "Group A"
        group_b = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(1), Inches(.5)); group_b.text = "Group B"
        ungroup_a = slide.shapes.add_textbox(Inches(7), Inches(3), Inches(1), Inches(.5)); ungroup_a.text = "Ungroup A"
        ungroup_b = slide.shapes.add_textbox(Inches(7), Inches(4), Inches(1), Inches(.5)); ungroup_b.text = "Ungroup B"
        existing_group = slide.shapes.add_group_shape([ungroup_a, ungroup_b])
        source_pptx = root / "real.pptx"
        presentation.save(source_pptx)
        edited_pptx = root / "edited-powerpoint.pptx"
        powerpoint = edit(source_pptx, edited_pptx, "powerpoint", [
            {"type": "setShapeText", "slideIndex": 0, "shapeId": str(box.shape_id), "text": "Edited in OneDrive"},
            {"type": "setTextStyle", "slideIndex": 0, "shapeId": str(box.shape_id), "bold": True, "fontSize": 24, "color": "2563EB"},
            {"type": "addTextBox", "slideIndex": 0, "text": "Native text box", "x": 914400, "y": 4114800, "width": 2743200, "height": 457200},
            {"type": "deleteShape", "slideIndex": 0, "shapeId": str(disposable.shape_id)},
            {"type": "replaceImage", "slideIndex": 0, "shapeId": str(picture.shape_id), "base64": base64.b64encode(source_png).decode("ascii"), "contentType": "image/png"},
            {"type": "duplicateSlide", "slideIndex": 0},
        ])
        rich_pptx = root / "rich-powerpoint.pptx"
        rich_powerpoint_operations = [
            {"type": "addSlide", "afterIndex": 0},
            {"type": "addImage", "slideIndex": 0, "base64": image_b64, "contentType": "image/png", "x": 457200, "y": 4572000, "width": 457200, "height": 457200, "altText": "Added image"},
            {"type": "cropImage", "slideIndex": 0, "shapeId": str(picture.shape_id), "left": 0.1, "right": 0.1},
            {"type": "addTable", "slideIndex": 0, "rows": [["One", "Two"], ["Three", "Four"]], "x": 3657600, "y": 4572000, "width": 1828800, "height": 914400},
            {"type": "insertTableRow", "slideIndex": 0, "shapeId": str(table_shape.shape_id), "rowIndex": 1, "values": ["R1", "R2"]},
            {"type": "insertTableColumn", "slideIndex": 0, "shapeId": str(table_shape.shape_id), "columnIndex": 1, "values": ["C1", "C2", "C3"]},
            {"type": "deleteTableRow", "slideIndex": 0, "shapeId": str(table_shape.shape_id), "rowIndex": 2},
            {"type": "deleteTableColumn", "slideIndex": 0, "shapeId": str(table_shape.shape_id), "columnIndex": 2},
            {"type": "setShapeAltText", "slideIndex": 0, "shapeId": str(box.shape_id), "title": "Main box", "description": "Edited content"},
            {"type": "setZOrder", "slideIndex": 0, "shapeId": str(box.shape_id), "position": "front"},
            {"type": "groupShapes", "slideIndex": 0, "shapeIds": [str(group_a.shape_id), str(group_b.shape_id)], "name": "Codex Group"},
            {"type": "ungroupShape", "slideIndex": 0, "shapeId": str(existing_group.shape_id)},
            {"type": "applySlideLayout", "slideIndex": 1, "layoutName": presentation.slide_layouts[6].name},
        ]
        rich_powerpoint = edit(source_pptx, rich_pptx, "powerpoint", rich_powerpoint_operations)
        mixed_powerpoint_path = root / "mixed-order-powerpoint.pptx"
        mixed_powerpoint = edit(source_pptx, mixed_powerpoint_path, "powerpoint", [
            {"type": "addSlide", "afterIndex": 0},
            {"type": "addTextBox", "slideIndex": 1, "shapeId": 50, "text": "Interleaved", "x": 914400, "y": 914400, "width": 1828800, "height": 457200},
            {"type": "setShapeAltText", "slideIndex": 1, "shapeId": "50", "title": "Ordered", "description": "Rich after legacy"},
        ])

        Document(edited_docx)
        reopened_excel = load_workbook(edited_xlsx, data_only=False)
        Presentation(edited_pptx)
        Document(rich_word_docx); load_workbook(rich_xlsx, data_only=False); Presentation(rich_pptx)
        mixed_workbook = load_workbook(mixed_excel_path, data_only=False)
        mixed_presentation = Presentation(mixed_powerpoint_path)
        rendered = {}
        for path in (edited_docx, edited_xlsx, edited_pptx, rich_word_docx, rich_xlsx, rich_pptx):
            profile = root / ("profile-" + path.stem)
            conversion = subprocess.run([str(SOFFICE), "--headless", f"-env:UserInstallation={profile.as_uri()}", "--convert-to", "pdf", "--outdir", str(output), str(path)], cwd=root, capture_output=True, text=True, check=False)
            pdf = output / (path.stem + ".pdf")
            diagnostic = (conversion.stdout + "\n" + conversion.stderr).lower()
            rendered[path.stem] = {"exitCode": conversion.returncode, "bytes": pdf.stat().st_size if pdf.exists() else 0, "cleanOpen": conversion.returncode == 0 and not any(term in diagnostic for term in ("repair", "corrupt", "fatal error")), "stderr": conversion.stderr[-500:]}

        checks["wordRealPackage"] = word["changeCount"] == 2 and rendered["edited-word"]["bytes"] > 0 and rendered["edited-word"]["cleanOpen"]
        reopened_sheet = reopened_excel["Data"]
        checks["excelRealPackage"] = (
            excel["changeCount"] == 5
            and reopened_sheet.tables["RevenueTable"].ref == "A1:B4"
            and len(reopened_sheet._charts) == 2
            and reopened_sheet["A2"].value == "South"
            and reopened_sheet["A3"].value == "West"
            and reopened_sheet["B4"].value == "=SUBTOTAL(109,[Revenue])"
            and rendered["edited-excel"]["bytes"] > 0
            and rendered["edited-excel"]["cleanOpen"]
        )
        checks["powerpointRealPackage"] = powerpoint["changeCount"] == 6 and rendered["edited-powerpoint"]["bytes"] > 0 and rendered["edited-powerpoint"]["cleanOpen"]
        rich_coverage = {"word": {entry["type"] for entry in rich_word_operations}, "excel": {entry["type"] for entry in rich_excel_operations}, "powerpoint": {entry["type"] for entry in rich_powerpoint_operations}}
        checks["richOperationCoverage"] = rich_coverage == EXPECTED_RICH_OPERATIONS
        checks["wordRichPackage"] = rich_word["changeCount"] == len(rich_word_operations) and rendered["rich-word"]["cleanOpen"]
        checks["excelRichPackage"] = rich_excel["changeCount"] == len(rich_excel_operations) and rendered["rich-excel"]["cleanOpen"]
        checks["powerpointRichPackage"] = rich_powerpoint["changeCount"] == len(rich_powerpoint_operations) and rendered["rich-powerpoint"]["cleanOpen"]
        checks["mixedExcelOperationOrder"] = mixed_excel["changeCount"] == 3 and mixed_workbook["Interleaved"]["A1"].value == "ordered" and mixed_workbook["Interleaved"]["A2"].hyperlink is not None
        mixed_shape = next(shape for shape in mixed_presentation.slides[1].shapes if str(shape.shape_id) == "50")
        mixed_c_nv_pr = mixed_shape._element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
        checks["mixedPowerPointOperationOrder"] = mixed_powerpoint["changeCount"] == 3 and mixed_shape.text == "Interleaved" and mixed_c_nv_pr.get("title") == "Ordered" and mixed_c_nv_pr.get("descr") == "Rich after legacy"
        checks["consolidatedFixturesAreGenuinePackages"] = all(entry["bytes"] > 0 and entry["cleanOpen"] for entry in emitted_rendered.values())
        print(json.dumps({"ok": all(checks.values()), "checks": checks, "rendered": rendered, "consolidatedRendered": emitted_rendered}, indent=2))
        raise SystemExit(0 if all(checks.values()) else 1)


if __name__ == "__main__":
    main()
